import collections 
import collections.abc
from pptx import Presentation
import sys

try:
    prs = Presentation("Société BARRANI Bâtiment - Portfolio Overview.pptx")
    for i, slide in enumerate(prs.slides):
        print(f"--- Slide {i+1} ---")
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    print(run.text)
except Exception as e:
    print(f"Error: {e}")
